from flask import Blueprint, request, jsonify
from db import execute_query, execute_write
from rag_service import add_document
import uuid
import os
import yt_dlp
import speech_recognition as sr
import imageio_ffmpeg
import subprocess

knowledge_bp = Blueprint('knowledge_bp', __name__)

@knowledge_bp.route('/upload', methods=['POST'])
def upload_document():
    files = request.files.getlist('files')
    if not files and 'file' in request.files:
        files = [request.files['file']]
        
    if not files:
        return jsonify({"success": False, "message": "No file part"}), 400
        
    plan_id = request.form.get('plan_id')
    kt_day = request.form.get('kt_day')
    
    if not plan_id:
        return jsonify({"success": False, "message": "Missing plan_id"}), 400
        
    try:
        plan_id = int(plan_id)
    except ValueError:
        return jsonify({"success": False, "message": "Invalid plan_id"}), 400
        
    manager_id = -1
    auth_header = request.headers.get('Authorization')
    if auth_header and auth_header.startswith('Bearer '):
        token = auth_header.split(' ')[1]
        try:
            import jwt
            from config import Config
            payload = jwt.decode(token, Config.SECRET_KEY, algorithms=["HS256"])
            user_email = payload.get('email')
            user_role = payload.get('role')
            user_id = payload.get('sub')
            
            user_full_name = None
            if user_id:
                users = execute_query("SELECT full_name FROM users WHERE id = %s", (user_id,))
                if users:
                    user_full_name = users[0]['full_name']
            
            from services.plan_service import resolve_stakeholder_for_user
            if user_email:
                manager_id = resolve_stakeholder_for_user(user_email, user_full_name, user_role)
        except Exception:
            pass

    import datetime
    timestamp = datetime.datetime.utcnow().isoformat()
    results = []

    try:
        for file in files:
            if file.filename == '':
                continue
                
            filename = file.filename
            
            if filename == 'CONFLUENCE_SYNC.txt':
                from connectors import ConfluenceConnector
                confluence = ConfluenceConnector()
                kb_chunks = confluence.fetch_kb_from_confluence()
                
                text = "\n".join(kb_chunks)
                filename = 'confluence_auto_sync'
            else:
                ext = os.path.splitext(filename)[1].lower()
                text = ""
                
                if ext == '.pdf':
                    import pypdf
                    pdf_reader = pypdf.PdfReader(file)
                    for page in pdf_reader.pages:
                        text += page.extract_text() + "\n"
                elif ext in ['.docx', '.doc', '.docs']:
                    import docx
                    doc = docx.Document(file)
                    for para in doc.paragraphs:
                        text += para.text + "\n"
                elif ext in ['.ppt', '.pptx']:
                    from pptx import Presentation
                    prs = Presentation(file)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                elif ext == '.txt':
                    text = file.read().decode('utf-8', errors='ignore')
                else:
                    continue # Skip unsupported files in multi-upload
                
            doc_id = str(uuid.uuid4())
            metadata = {
                "plan_id": plan_id,
                "day": kt_day,
                "manager_id": manager_id,
                "file_name": filename,
                "created_at": timestamp
            }
            
            chunk_count = add_document(doc_id, text, metadata)
            
            query = """
                INSERT INTO knowledge_documents (plan_id, kt_day, filename, chunk_count)
                VALUES (%s, %s, %s, %s)
            """
            doc_db_id = execute_write(query, (plan_id, kt_day, filename, chunk_count))
            
            results.append({
                "id": doc_db_id,
                "plan_id": plan_id,
                "kt_day": kt_day,
                "filename": filename,
                "chunk_count": chunk_count
            })
            
        return jsonify({
            "success": True, 
            "data": results,
            "message": f"Processed {len(results)} document(s) successfully"
        }), 201
        
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500

@knowledge_bp.route('/plan/<int:plan_id>', methods=['GET'])
def get_plan_documents(plan_id):
    try:
        query = "SELECT * FROM knowledge_documents WHERE plan_id = %s ORDER BY uploaded_at DESC"
        docs = execute_query(query, (plan_id,))
        return jsonify({"success": True, "data": docs}), 200
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500

def transcribe_audio_from_url(video_url, user_id=None):
    import re
    import requests
    from db import execute_query
    
    ffmpeg_path = imageio_ffmpeg.get_ffmpeg_exe()
    
    uid = str(uuid.uuid4())
    temp_folder = os.path.join(os.getcwd(), 'temp_audio_files')
    os.makedirs(temp_folder, exist_ok=True)
    
    is_google_drive = "drive.google.com" in video_url
    is_onedrive = "sharepoint.com" in video_url or "1drv.ms" in video_url or "onedrive.live.com" in video_url
    
    temp_orig = None
    ext = 'mp4' # Default fallback for API downloads
    
    if is_google_drive and user_id:
        try:
            file_id_match = re.search(r'/d/([^/]+)', video_url) or re.search(r'id=([^&]+)', video_url)
            if not file_id_match:
                return "Failed to extract audio. Could not parse Google Drive File ID."
            file_id = file_id_match.group(1)
            
            users = execute_query("SELECT google_token FROM users WHERE id = %s", (user_id,))
            if not users or not users[0].get('google_token'):
                return "Failed to extract audio. Please connect your Google Drive account first."
                
            import json
            from google.oauth2.credentials import Credentials
            from googleapiclient.discovery import build
            
            token_data = json.loads(users[0]['google_token'])
            creds = Credentials.from_authorized_user_info(token_data)
            service = build('drive', 'v3', credentials=creds)
            
            request = service.files().get_media(fileId=file_id)
            import io
            from googleapiclient.http import MediaIoBaseDownload
            temp_orig = os.path.join(temp_folder, f"temp_audio_{uid}.{ext}")
            with io.FileIO(temp_orig, 'wb') as fh:
                downloader = MediaIoBaseDownload(fh, request)
                done = False
                while done is False:
                    status, done = downloader.next_chunk()
        except Exception as e:
            if "insufficient permissions" in str(e).lower() or "forbidden" in str(e).lower() or "invalid_grant" in str(e).lower():
                return "Failed to extract audio. Please reconnect your Google Drive account."
            return f"Failed to extract audio from Google Drive: {e}"

    elif is_onedrive and user_id:
        try:
            users = execute_query("SELECT ms_token FROM users WHERE id = %s", (user_id,))
            if not users or not users[0].get('ms_token'):
                return "Failed to extract audio. Please connect your Microsoft account first."
                
            import json
            token_data = json.loads(users[0]['ms_token'])
            access_token = token_data.get('access_token')
            
            import base64
            encoded_url = base64.urlsafe_b64encode(video_url.encode()).decode().rstrip('=')
            share_id = f"u!{encoded_url}"
            graph_url = f"https://graph.microsoft.com/v1.0/shares/{share_id}/driveItem/content"
            
            headers = {'Authorization': f'Bearer {access_token}'}
            resp = requests.get(graph_url, headers=headers, stream=True)
            if resp.status_code == 200:
                temp_orig = os.path.join(temp_folder, f"temp_audio_{uid}.{ext}")
                with open(temp_orig, 'wb') as f:
                    for chunk in resp.iter_content(chunk_size=8192):
                        f.write(chunk)
            elif resp.status_code == 401 or resp.status_code == 403:
                return "Failed to extract audio. Please reconnect your Microsoft account."
            else:
                return f"Failed to extract audio from OneDrive: HTTP {resp.status_code}"
        except Exception as e:
            return f"Failed to extract audio from OneDrive: {e}"

    else:
        ydl_opts = {
            'format': 'bestaudio/best',
            'quiet': True,
            'noplaylist': True
        }
        ydl_opts['outtmpl'] = os.path.join(temp_folder, f'temp_audio_{uid}.%(ext)s')
        
        try:
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                info = ydl.extract_info(video_url, download=True)
                ext = info.get('ext', 'm4a')
        except Exception as e:
            return f"Failed to extract audio from URL. Please ensure it is publicly accessible. Error: {e}"
            
        temp_orig = os.path.join(temp_folder, f"temp_audio_{uid}.{ext}")
        
    temp_wav = os.path.join(temp_folder, f"temp_audio_{uid}.wav")
    
    try:
        subprocess.run([ffmpeg_path, "-y", "-i", temp_orig, "-ac", "1", "-ar", "16000", temp_wav], check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    except Exception as e:
        if temp_orig and os.path.exists(temp_orig): os.remove(temp_orig)
        return f"Conversion failed: {e}"
        
    r = sr.Recognizer()
    full_text = ""
    try:
        with sr.AudioFile(temp_wav) as source:
            while True:
                audio_data = r.record(source, duration=55) # 55 seconds chunk
                if not audio_data.frame_data:
                    break
                try:
                    text = r.recognize_google(audio_data)
                    full_text += text + " "
                except sr.UnknownValueError:
                    pass
                except sr.RequestError as e:
                    full_text += f"[API Error: {e}] "
    except Exception as e:
        full_text = f"Transcription failed: {e}"
        
    if temp_orig and os.path.exists(temp_orig): os.remove(temp_orig)
    if os.path.exists(temp_wav): os.remove(temp_wav)
    
    return full_text.strip()

@knowledge_bp.route('/extract-transcript', methods=['POST'])
def extract_transcript():
    data = request.json
    if not data or not data.get('url'):
        return jsonify({"success": False, "message": "Missing video URL"}), 400
        
    url = data['url']
    
    user_id = None
    auth_header = request.headers.get('Authorization')
    if auth_header and auth_header.startswith('Bearer '):
        token = auth_header.split(' ')[1]
        try:
            import jwt
            from config import Config
            payload = jwt.decode(token, Config.SECRET_KEY, algorithms=["HS256"])
            user_id = payload.get('sub')
        except Exception:
            pass
    
    try:
        transcript_text = transcribe_audio_from_url(url, user_id)
        if not transcript_text:
            transcript_text = "Could not transcribe audio (no speech detected)."
            
        if transcript_text.startswith("Failed to extract audio") or transcript_text.startswith("Conversion failed") or transcript_text.startswith("Transcription failed"):
            return jsonify({"success": False, "message": transcript_text}), 400
            
        return jsonify({
            "success": True, 
            "data": {
                "transcript": transcript_text
            },
            "message": "Transcript extracted successfully"
        }), 200
    except Exception as e:
        return jsonify({"success": False, "message": f"Transcription error: {str(e)}"}), 500

@knowledge_bp.route('/upload-transcript', methods=['POST'])
def upload_transcript():
    data = request.json
    if not data:
        return jsonify({"success": False, "message": "Missing JSON body"}), 400
        
    plan_id = data.get('plan_id')
    kt_day = data.get('kt_day')
    text = data.get('text')
    url = data.get('url', 'Unknown URL')
    
    if not plan_id or not text:
        return jsonify({"success": False, "message": "Missing plan_id or text"}), 400
        
    try:
        plan_id = int(plan_id)
    except ValueError:
        return jsonify({"success": False, "message": "Invalid plan_id"}), 400
        
    try:
        filename = f"Transcript_Day_{kt_day}.txt" if kt_day else f"Transcript_URL.txt"
        
        doc_id = str(uuid.uuid4())
        metadata = {"plan_id": plan_id, "filename": filename, "kt_day": kt_day, "source_url": url}
        
        chunk_count = add_document(doc_id, text, metadata)
        
        query = """
            INSERT INTO knowledge_documents (plan_id, kt_day, filename, chunk_count)
            VALUES (%s, %s, %s, %s)
        """
        doc_db_id = execute_write(query, (plan_id, kt_day, filename, chunk_count))
        
        return jsonify({
            "success": True, 
            "data": {
                "id": doc_db_id,
                "plan_id": plan_id,
                "kt_day": kt_day,
                "filename": filename,
                "chunk_count": chunk_count
            },
            "message": "Transcript processed and added to knowledge base"
        }), 201
        
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500

