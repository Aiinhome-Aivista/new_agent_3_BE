from flask import Blueprint, request, jsonify, send_file
from db import execute_query, execute_write
from llm_service import call_llm
import os
from datetime import datetime

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from docx import Document
except ImportError:
    Document = None

reporting_bp = Blueprint('reporting_bp', __name__)

REPORTS_DIR = os.path.join(os.path.dirname(__file__), '..', 'reports_output')

@reporting_bp.route('/weekly', methods=['POST'])
def generate_weekly():
    data = request.json
    if 'plan_id' not in data:
        return jsonify({"success": False, "message": "Missing plan_id"}), 400
        
    plan_id = data['plan_id']
    try:
        from services.reporting_service import generate_weekly_service
        result_data = generate_weekly_service(plan_id)
        return jsonify({"success": True, "message": "Weekly report generated", "data": result_data}), 201
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500

@reporting_bp.route('/final', methods=['POST'])
def generate_final():
    data = request.json
    if 'plan_id' not in data:
        return jsonify({"success": False, "message": "Missing plan_id"}), 400
        
    plan_id = data['plan_id']
    try:
        from services.reporting_service import generate_final_service
        result_data = generate_final_service(plan_id)
        return jsonify({"success": True, "message": "Final report generated", "data": result_data}), 201
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500

@reporting_bp.route('/', methods=['GET'])
def get_reports():
    try:
        query = "SELECT * FROM reports ORDER BY generated_at DESC"
        reports = execute_query(query)
        # Ensure filenames ending with .docx are displayed as .pptx in file_path
        for r in reports:
            if r.get('file_path') and r['file_path'].endswith('.docx'):
                r['file_path'] = r['file_path'][:-5] + '.pptx'
        return jsonify({"success": True, "data": reports}), 200
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500

@reporting_bp.route('/download/<int:id>', methods=['GET'])
def download_report(id):
    try:
        query = "SELECT file_path FROM reports WHERE id = %s"
        report = execute_query(query, (id,))
        if not report:
            return jsonify({"success": False, "message": "Report not found"}), 404
            
        raw_filename = report[0]['file_path']
        filepath = os.path.join(REPORTS_DIR, raw_filename)
        
        # If .pptx exists, send it directly
        if os.path.exists(filepath) and raw_filename.endswith('.pptx'):
            return send_file(filepath, as_attachment=True, download_name=raw_filename, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')

        # Check if corresponding .pptx or original .docx exists
        pptx_filename = raw_filename if raw_filename.endswith('.pptx') else raw_filename.rsplit('.', 1)[0] + '.pptx'
        pptx_filepath = os.path.join(REPORTS_DIR, pptx_filename)

        if os.path.exists(pptx_filepath):
            return send_file(pptx_filepath, as_attachment=True, download_name=pptx_filename, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')

        # If docx exists, convert on the fly to pptx
        docx_filepath = os.path.join(REPORTS_DIR, raw_filename if raw_filename.endswith('.docx') else raw_filename.rsplit('.', 1)[0] + '.docx')
        if os.path.exists(docx_filepath):
            from services.reporting_service import generate_report_pptx
            doc = Document(docx_filepath) if Document else None
            text_lines = [p.text for p in doc.paragraphs] if doc else []
            content_str = "\n".join(text_lines)
            generate_report_pptx(f"KT Status Report", content_str, pptx_filename)
            return send_file(pptx_filepath, as_attachment=True, download_name=pptx_filename, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')

        return jsonify({"success": False, "message": "File not found on disk"}), 404
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500

@reporting_bp.route('/view/<int:id>', methods=['GET'])
def view_report(id):
    try:
        query = "SELECT file_path FROM reports WHERE id = %s"
        report = execute_query(query, (id,))
        if not report:
            return jsonify({"success": False, "message": "Report not found"}), 404
            
        raw_filename = report[0]['file_path']
        pptx_filename = raw_filename if raw_filename.endswith('.pptx') else raw_filename.rsplit('.', 1)[0] + '.pptx'
        filepath = os.path.join(REPORTS_DIR, raw_filename)
        pptx_filepath = os.path.join(REPORTS_DIR, pptx_filename)
        
        target_path = pptx_filepath if os.path.exists(pptx_filepath) else filepath
        
        if not os.path.exists(target_path):
            # Check if docx exists to parse as PPT slides
            docx_filepath = os.path.join(REPORTS_DIR, raw_filename if raw_filename.endswith('.docx') else raw_filename.rsplit('.', 1)[0] + '.docx')
            if os.path.exists(docx_filepath) and Document:
                doc = Document(docx_filepath)
                paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
                slides_data = [{
                    "slide_number": 1,
                    "title": paragraphs[0] if paragraphs else "Report Overview",
                    "subtitle": "Knowledge Transfer Management Summary",
                    "content": paragraphs[1:] if len(paragraphs) > 1 else []
                }]
                return jsonify({
                    "success": True,
                    "data": {
                        "is_ppt": True,
                        "filename": pptx_filename,
                        "total_slides": len(slides_data),
                        "slides": slides_data
                    }
                }), 200
            return jsonify({"success": False, "message": "File not found on disk"}), 404
            
        if not Presentation:
            return jsonify({"success": False, "message": "python-pptx library not available"}), 500
            
        import re
        def clean_markdown(text):
            if not text:
                return ""
            text = text.replace('**', '').strip()
            text = re.sub(r'^\s*[-*•]\s*', '', text)
            text = re.sub(r'^\s*#{1,6}\s*', '', text)
            text = re.sub(r'^\s*\d+\.\s*', '', text)
            return text.strip()

        prs = Presentation(target_path)
        slides_data = []
        for idx, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        txt = clean_markdown(p.text)
                        if txt:
                            texts.append(txt)
            
            title = texts[0] if texts else f"Slide {idx + 1}"
            subtitle = texts[1] if len(texts) > 1 and len(texts[1]) < 90 and not texts[1].startswith('•') else ""
            content_items = [clean_markdown(item) for item in (texts[2:] if subtitle else texts[1:]) if clean_markdown(item)]
            
            slides_data.append({
                "slide_number": idx + 1,
                "title": title,
                "subtitle": subtitle,
                "content": content_items
            })
            
        return jsonify({
            "success": True,
            "data": {
                "is_ppt": True,
                "filename": pptx_filename,
                "total_slides": len(slides_data),
                "slides": slides_data
            }
        }), 200
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500

@reporting_bp.route('/<int:id>/status', methods=['PUT', 'PATCH'])
def update_report_status(id):
    data = request.json or {}
    new_status = data.get('status', 'approved')
    try:
        query = "UPDATE reports SET status = %s WHERE id = %s"
        execute_write(query, (new_status, id))
        return jsonify({"success": True, "message": f"Report status updated to {new_status}"}), 200
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500

