import os
import json
import logging
from llm_service import call_llm

def extract_holiday_info_from_doc_service(uploaded_files):
    """
    Extracts text from uploaded documents (PDF, DOC/DOCX, PPT/PPTX, TXT)
    and prompts the LLM to identify holiday dates, names, and years.
    Returns a list of dictionaries with keys: holiday_date, holiday_name, holiday_year.
    """
    combined_text = ""
    for file_storage in uploaded_files:
        filename = file_storage.filename.lower()
        ext = os.path.splitext(filename)[1]
        doc_text = ""
        file_storage.seek(0)

        if ext == '.pdf':
            try:
                import pypdf
                pdf_reader = pypdf.PdfReader(file_storage)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        doc_text += page_text + "\n"
            except Exception as e:
                logging.error(f"Error reading PDF {filename}: {e}")
                raise Exception(f"Failed to parse PDF file ({filename}): {str(e)}")
        elif ext in ['.docx', '.doc', '.docs']:
            try:
                import docx
                doc = docx.Document(file_storage)
                for para in doc.paragraphs:
                    if para.text and para.text.strip():
                        doc_text += para.text.strip() + "\n"
                for table in doc.tables:
                    for row in table.rows:
                        row_cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
                        if row_cells:
                            doc_text += " | ".join(row_cells) + "\n"
            except Exception as e:
                logging.error(f"Error reading DOC/DOCX {filename}: {e}")
                raise Exception(f"Failed to parse Word document ({filename}): {str(e)}")
        elif ext in ['.pptx', '.ppt']:
            try:
                from pptx import Presentation
                prs = Presentation(file_storage)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            doc_text += shape.text + "\n"
            except Exception as e:
                logging.error(f"Error reading PPT/PPTX {filename}: {e}")
                raise Exception(f"Failed to parse PowerPoint document ({filename}): {str(e)}")
        elif ext == '.txt':
            doc_text = file_storage.read().decode('utf-8', errors='ignore')
        else:
            raise Exception(f"Unsupported file type for {filename}. Allowed: .pdf, .doc, .docx, .ppt, .pptx, .txt")

        doc_text = doc_text.strip()
        if doc_text:
            combined_text += f"\n--- Document: {filename} ---\n{doc_text}\n"

    combined_text = combined_text.strip()
    if not combined_text:
        raise Exception("The uploaded document(s) are empty or text could not be extracted.")

    import re
    from datetime import datetime
    try:
        from dateutil import parser
    except ImportError:
        raise Exception("python-dateutil library is required for date parsing.")

    all_holidays = []
    
    # 1. Dynamic Year Extraction (Fallback to current year)
    current_year = datetime.now().year
    year_match = re.search(r'\b(202[0-9])\b', combined_text)
    if year_match:
        current_year = int(year_match.group(1))

    # Regex patterns for cleaning holiday names
    junk_patterns = [
        r'\bYes\b', r'\bNo\b', r'\bTotal\s*Holidays?\b', r'\bIndia\s*Holiday\s*List\b',
        r'\bDun\s*&\s*Bradstreet\b', r'\bRestricted\s*Sensitive\b',
        r'^\d{1,2}[\.\)]?\s+' # leading serial numbers like "1." or "1)"
    ]

    months = r'(Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*'
    
    date_regex = re.compile(
        rf'(\b\d{{1,4}}[-/]\d{{1,2}}[-/]\d{{1,4}}\b)|'
        rf'(\b{months}\s+\d{{1,2}}(?:st|nd|rd|th)?(?:,?\s+\d{{4}})?\b)|'
        rf'(\b\d{{1,2}}(?:st|nd|rd|th)?\s+{months}(?:,?\s+\d{{4}})?\b)', 
        re.IGNORECASE
    )

    lines = combined_text.split('\n')
    for line in lines:
        line = line.strip()
        if not line or len(line) < 5:
            continue
            
        try:
            match = date_regex.search(line)
            if not match:
                continue
                
            date_str = match.group(0)
            
            # Remove the exact matched date string
            name_str = line.replace(date_str, '')
            
            # Remove any trailing or leading day of week
            days = r'\b(Monday|Tuesday|Wednesday|Thursday|Friday|Saturday|Sunday|Mon|Tue|Wed|Thu|Fri|Sat|Sun)\b,?'
            name_str = re.sub(days, '', name_str, flags=re.IGNORECASE)
            
            # Parse the extracted date string
            parsed_date = parser.parse(date_str, default=datetime(current_year, 1, 1))
            
            # Clean up the name
            name = name_str
            for pattern in junk_patterns:
                name = re.sub(pattern, '', name, flags=re.IGNORECASE)
            
            name = name.strip()
            # Trim punctuation from edges
            name = re.sub(r'^[-|:\s,\.]+|[-|:\s,\.]+$', '', name).strip()
            
            if name and len(name) > 2 and not name.isdigit():
                # Avoid appending header lines that accidentally matched a year
                if "holiday list" not in name.lower() and "holiday calendar" not in name.lower():
                    all_holidays.append({
                        "date": parsed_date.strftime("%Y-%m-%d"),
                        "name": name,
                        "year": parsed_date.year
                    })
        except Exception:
            pass

    if not all_holidays:
        raise Exception("Failed to extract any holidays from the document using direct parsing.")

    # Remove exact duplicates
    unique_holidays = []
    seen = set()
    for h in all_holidays:
        identifier = f"{h.get('date')}-{h.get('name')}"
        if identifier not in seen:
            seen.add(identifier)
            unique_holidays.append(h)

    return unique_holidays
