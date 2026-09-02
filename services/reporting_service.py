from db import execute_query, execute_write
from llm_service import call_llm, load_prompt
import os
from datetime import datetime
import re

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    Presentation = None

REPORTS_DIR = os.path.join(os.path.dirname(__file__), '..', 'reports_output')
os.makedirs(REPORTS_DIR, exist_ok=True)

def clean_markdown(text):
    if not text:
        return ""
    text = text.replace('**', '').strip()
    text = re.sub(r'^\s*[-*•]\s*', '', text)    # Strip leading bullet icons
    text = re.sub(r'^\s*#{1,6}\s*', '', text)  # Strip leading hashes (#, ##, ###, ####)
    text = re.sub(r'^\s*\d+\.\s*', '', text)    # Strip leading numbers (1., 2.)
    return text.strip()

def generate_report_pptx(title, content, filename):
    if not Presentation:
        raise Exception("python-pptx is not installed")
        
    prs = Presentation()
    prs.slide_width = Inches(13.333) # 16:9 Widescreen
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    title_slide = prs.slides.add_slide(blank_layout)

    # Top accent bar
    top_bar = title_slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.35))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = RGBColor(217, 83, 30) # PwC Orange
    top_bar.line.fill.background()

    # Cover Slide Header
    tb = title_slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = clean_markdown(title)
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 41, 59)

    p_sub = tf.add_paragraph()
    p_sub.text = f"Generated on {datetime.now().strftime('%B %d, %Y')} | Delivery & Solution Advisory"
    p_sub.font.size = Pt(16)
    p_sub.font.color.rgb = RGBColor(217, 83, 30)
    p_sub.space_before = Pt(10)

    # Cover Slide Metadata & Overview Box
    tb_meta = title_slide.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(11.333), Inches(2.5))
    tf_meta = tb_meta.text_frame
    tf_meta.word_wrap = True
    p_m1 = tf_meta.paragraphs[0]
    p_m1.text = "Executive Program Status Report"
    p_m1.font.size = Pt(18)
    p_m1.font.bold = True
    p_m1.font.color.rgb = RGBColor(30, 41, 59)

    p_m2 = tf_meta.add_paragraph()
    p_m2.text = "• Target Persona: Delivery / Engagement Manager & PwC Leadership Review"
    p_m2.font.size = Pt(14)
    p_m2.font.color.rgb = RGBColor(71, 85, 105)
    p_m2.space_before = Pt(8)

    p_m3 = tf_meta.add_paragraph()
    p_m3.text = "• Governance & Verification: RAG-Grounded & Audit Verified Status Update"
    p_m3.font.size = Pt(14)
    p_m3.font.color.rgb = RGBColor(71, 85, 105)
    p_m3.space_before = Pt(6)

    # Process markdown content into slide cards
    lines = [l.strip() for l in content.split('\n') if l.strip()]
    current_slide_title = "Executive Status Overview"
    current_bullets = []

    def flush_slide(stitle, bullets):
        if not bullets and not stitle:
            return
        s = prs.slides.add_slide(blank_layout)
        
        # Header accent bar
        bar = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(217, 83, 30)
        bar.line.fill.background()

        # Slide Title
        tb_t = s.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
        tf_t = tb_t.text_frame
        tf_t.word_wrap = True
        pt = tf_t.paragraphs[0]
        pt.text = clean_markdown(stitle)
        pt.font.size = Pt(22)
        pt.font.bold = True
        pt.font.color.rgb = RGBColor(217, 83, 30)

        # Slide Body Bullets
        tb_b = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3))
        tf_b = tb_b.text_frame
        tf_b.word_wrap = True

        valid_bullets = [clean_markdown(b) for b in bullets if clean_markdown(b)]
        for idx, bullet in enumerate(valid_bullets[:10]):
            pb = tf_b.paragraphs[0] if idx == 0 else tf_b.add_paragraph()
            pb.text = f"•  {bullet}"
            pb.font.size = Pt(14)
            pb.font.color.rgb = RGBColor(51, 65, 85)
            pb.space_after = Pt(10)

    for line in lines:
        cleaned_line = clean_markdown(line)
        if not cleaned_line:
            continue
            
        if line.startswith('#') or (line.endswith(':') and len(cleaned_line) < 45):
            if current_bullets:
                flush_slide(current_slide_title, current_bullets)
                current_bullets = []
            current_slide_title = cleaned_line
        else:
            current_bullets.append(cleaned_line)

    if current_bullets:
        flush_slide(current_slide_title, current_bullets)

    filepath = os.path.join(REPORTS_DIR, filename)
    prs.save(filepath)
    return filepath

def generate_weekly_service(plan_id):
    comp_query = """
        SELECT 
            (SELECT COUNT(*) FROM completion_tracking WHERE plan_id = %s AND completion_percent = 100) as completed_topics,
            (SELECT COUNT(*) FROM plan_topics WHERE plan_id = %s) as total_topics
        FROM DUAL
    """
    comp_res = execute_query(comp_query, (plan_id, plan_id))
    
    if comp_res and comp_res[0]['total_topics'] and int(comp_res[0]['total_topics']) > 0:
        avg_comp = (float(comp_res[0]['completed_topics']) / float(comp_res[0]['total_topics'])) * 100.0
    else:
        avg_comp = 0.0
    
    risk_query = "SELECT description, severity FROM risks WHERE plan_id = %s AND status != 'resolved'"
    risks = execute_query(risk_query, (plan_id,))
    
    prompt = load_prompt("report_weekly_summary.txt", avg_comp=avg_comp, risks=risks)
    summary = call_llm(prompt)
    
    filename = f"Weekly_Report_{plan_id}_{datetime.now().strftime('%Y%m%d%H%M%S')}.pptx"
    filepath = generate_report_pptx(f"Weekly KT Report (Plan {plan_id})", summary, filename)
    
    query = "INSERT INTO reports (plan_id, report_type, file_path) VALUES (%s, %s, %s)"
    report_id = execute_write(query, (plan_id, 'weekly', filename))
    
    return {"id": report_id, "filename": filename}

def generate_final_service(plan_id):
    plan_query = "SELECT application_name FROM kt_plans WHERE id = %s"
    app_name = execute_query(plan_query, (plan_id,))[0]['application_name']
    
    topics_query = "SELECT topic FROM completion_tracking WHERE plan_id = %s AND completion_percent = 100"
    completed_topics_res = execute_query(topics_query, (plan_id,))
    completed_topics_list = [row['topic'] for row in completed_topics_res]
    source_type = "completed topics (100% progress)"
    
    if not completed_topics_list:
        fallback_query = "SELECT topic FROM completion_tracking WHERE plan_id = %s AND completion_percent > 0"
        completed_topics_res = execute_query(fallback_query, (plan_id,))
        completed_topics_list = [row['topic'] for row in completed_topics_res]
        source_type = "partially covered topics (progress > 0%)"
        
    if not completed_topics_list:
        fallback_query = "SELECT topic_name FROM plan_topics WHERE plan_id = %s"
        completed_topics_res = execute_query(fallback_query, (plan_id,))
        completed_topics_list = [row['topic_name'] for row in completed_topics_res]
        source_type = "all topics in the plan"
        
    topics_text = "\n".join([f"- {t}" for t in completed_topics_list]) if completed_topics_list else "No topics found"
    
    manager_query = "SELECT s.name FROM stakeholders s JOIN kt_plans kp ON kp.approved_by = s.id WHERE kp.id = %s"
    manager_res = execute_query(manager_query, (plan_id,))
    manager_name = manager_res[0]['name'] if manager_res else "[Manager Name Not Assigned]"
    
    giver_query = "SELECT DISTINCT s.name FROM stakeholders s JOIN attendance a ON a.stakeholder_id = s.id JOIN meetings m ON a.meeting_id = m.id WHERE m.plan_id = %s AND (s.role = 'Outgoing SME (Knowledge Giver)' OR s.role = 'outgoing_sme' OR s.role LIKE '%outgoing%' OR s.role LIKE '%Giver%')"
    giver_res = execute_query(giver_query, (plan_id,))
    giver_names = ", ".join([r['name'] for r in giver_res]) if giver_res else ""
    
    if not giver_names:
        global_giver = execute_query("SELECT name FROM stakeholders WHERE role = 'Outgoing SME (Knowledge Giver)' OR role = 'outgoing_sme' OR role LIKE '%outgoing%' OR role LIKE '%Giver%' LIMIT 1")
        giver_names = global_giver[0]['name'] if global_giver else "[Knowledge Giver]"
    
    receiver_query = "SELECT DISTINCT s.name FROM stakeholders s JOIN attendance a ON a.stakeholder_id = s.id JOIN meetings m ON a.meeting_id = m.id WHERE m.plan_id = %s AND (s.role = 'Incoming Team Member (Knowledge Receiver)' OR s.role = 'incoming_member' OR s.role LIKE '%incoming%' OR s.role LIKE '%Receiver%')"
    receiver_res = execute_query(receiver_query, (plan_id,))
    receiver_names = ", ".join([r['name'] for r in receiver_res]) if receiver_res else ""
    
    if not receiver_names:
        assess_query = "SELECT DISTINCT s.name FROM stakeholders s JOIN assessments a ON a.stakeholder_id = s.id WHERE a.plan_id = %s AND (s.role = 'Incoming Team Member (Knowledge Receiver)' OR s.role = 'incoming_member' OR s.role LIKE '%incoming%' OR s.role LIKE '%Receiver%')"
        assess_res = execute_query(assess_query, (plan_id,))
        receiver_names = ", ".join([r['name'] for r in assess_res]) if assess_res else ""
        
    if not receiver_names:
        global_receiver = execute_query("SELECT name FROM stakeholders WHERE role = 'Incoming Team Member (Knowledge Receiver)' OR role = 'incoming_member' OR role LIKE '%incoming%' OR role LIKE '%Receiver%' LIMIT 1")
        receiver_names = global_receiver[0]['name'] if global_receiver else "[Knowledge Receiver]"

    prompt = load_prompt(
        "report_final_assessment.txt",
        app_name=app_name,
        source_type=source_type,
        topics_text=topics_text,
        manager_name=manager_name,
        giver_names=giver_names,
        receiver_names=receiver_names
    )
    
    content = call_llm(prompt)
    
    filename = f"Final_Report_{plan_id}_{datetime.now().strftime('%Y%m%d%H%M%S')}.pptx"
    filepath = generate_report_pptx(f"Final KT Report - {app_name}", content, filename)
    
    query = "INSERT INTO reports (plan_id, report_type, file_path) VALUES (%s, %s, %s)"
    report_id = execute_write(query, (plan_id, 'final', filename))
    
    return {"id": report_id, "filename": filename}

